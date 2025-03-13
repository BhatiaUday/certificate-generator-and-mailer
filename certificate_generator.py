import pandas as pd
from pptx import Presentation
import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.base import MIMEBase
from email import encoders
import os
import logging
import json
from pathlib import Path
import re
import subprocess
import requests
import time

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[logging.FileHandler("certificate_generation.log"), logging.StreamHandler()]
)
logger = logging.getLogger(__name__)

def read_config(config_path='config.json'):
    with open(config_path, 'r') as f:
        return json.load(f)

def read_csv(csv_path):
    try:
        return pd.read_csv(csv_path)
    except Exception as e:
        logger.error(f"Error reading CSV: {e}")
        raise

def replace_placeholder_in_ppt(ppt_path, output_path, placeholder, replacement):
    """Replace placeholder in PPT with name and preserve text formatting"""
    try:
        presentation = Presentation(ppt_path)
        replacement_done = False
        
        for slide in presentation.slides:
            for shape in slide.shapes:
                if not hasattr(shape, "text_frame"):
                    continue
                    
                text_frame = shape.text_frame
                if placeholder not in text_frame.text:
                    continue
                    
                for paragraph in text_frame.paragraphs:
                    if placeholder not in paragraph.text:
                        continue
                        
                    para_alignment = paragraph.alignment
                    para_level = paragraph.level
                    para_space_before = paragraph.space_before
                    para_space_after = paragraph.space_after
                    para_line_spacing = paragraph.line_spacing
                    
                    runs_with_placeholder = []
                    for i, run in enumerate(paragraph.runs):
                        if placeholder in run.text:
                            runs_with_placeholder.append((i, run))
                    
                    if runs_with_placeholder:
                        if len(runs_with_placeholder) == 1:
                            idx, run = runs_with_placeholder[0]
                            
                            font = run.font
                            original_text = run.text
                            
                            run.text = original_text.replace(placeholder, replacement)
                        else:
                            logger.warning("Placeholder spans multiple runs - attempting to preserve formatting")
                            
                            full_text = paragraph.text
                            
                            new_text = full_text.replace(placeholder, replacement)
                            
                            for _ in range(len(paragraph.runs)):
                                paragraph._p.remove(paragraph.runs[0]._r)
                            
                            run = paragraph.add_run()
                            run.text = new_text
                            
                            _, source_run = runs_with_placeholder[0]
                            run.font.name = source_run.font.name
                            run.font.size = source_run.font.size
                            run.font.bold = source_run.font.bold
                            run.font.italic = source_run.font.italic
                            run.font.underline = source_run.font.underline
                            if hasattr(source_run.font, 'color') and source_run.font.color.rgb:
                                run.font.color.rgb = source_run.font.color.rgb
                        
                        paragraph.alignment = para_alignment
                        paragraph.level = para_level
                        paragraph.space_before = para_space_before
                        paragraph.space_after = para_space_after
                        paragraph.line_spacing = para_line_spacing
                        
                        replacement_done = True
        
        presentation.save(output_path)
        
        if not replacement_done:
            logger.warning(f"Placeholder '{placeholder}' not found or not replaceable in the presentation")
            
        return True
    except Exception as e:
        logger.error(f"Error processing PPT: {e}")
        return False

def convert_pptx_to_pdf_using_ilovepdf(pptx_path, pdf_path, api_key=None, api_secret=None):
    """Convert PPTX to PDF using ilovepdf.com API endpoints directly"""
    try:
        logger.info(f"Converting {pptx_path} to PDF using ilovepdf API")
        
        if not api_key:
            try:
                with open('config.json', 'r') as f:
                    config = json.load(f)
                    api_key = config.get("ilovepdf_public_key")
                    api_secret = config.get("ilovepdf_secret_key")
            except:
                api_key = os.environ.get("ILOVEPDF_PUBLIC_KEY")
                api_secret = os.environ.get("ILOVEPDF_SECRET_KEY")
        
        if not api_key or not api_secret:
            logger.error("API keys for ilovepdf not found. Add 'ilovepdf_public_key' and 'ilovepdf_secret_key' to config.json")
            return False
            
        base_url = "https://api.ilovepdf.com/v1"
        max_retries = 3
        retry_count = 0
        
        while retry_count < max_retries:
            try:
                auth_response = requests.post(
                    f"{base_url}/auth",
                    json={
                        "public_key": api_key,
                        "secret_key": api_secret
                    }
                )
                
                if auth_response.status_code != 200:
                    logger.error(f"Authentication failed: {auth_response.text}")
                    return False
                
                token = auth_response.json().get("token")
                
                response = requests.get(
                    f"{base_url}/start/officepdf",
                    headers={"Authorization": f"Bearer {token}"}
                )
                
                if response.status_code != 200:
                    logger.error(f"Failed to start task: {response.text}")
                    retry_count += 1
                    time.sleep(3)
                    continue
                
                task_data = response.json()
                task_id = task_data.get("task")
                server = task_data.get("server")
                
                with open(pptx_path, 'rb') as f:
                    files = {'file': (os.path.basename(pptx_path), f)}
                    response = requests.post(
                        f"https://{server}/v1/upload",
                        headers={"Authorization": f"Bearer {token}"},
                        data={"task": task_id},
                        files=files
                    )
                    
                if response.status_code != 200:
                    logger.error(f"Failed to upload file: {response.text}")
                    retry_count += 1
                    time.sleep(3)
                    continue
                    
                file_data = response.json()
                server_filename = file_data.get("server_filename")
                
                process_data = {
                    "task": task_id,
                    "tool": "officepdf",
                    "files": [{"server_filename": server_filename, "filename": os.path.basename(pptx_path)}]
                }
                
                response = requests.post(
                    f"https://{server}/v1/process",
                    headers={
                        "Authorization": f"Bearer {token}",
                        "Content-Type": "application/json"
                    },
                    json=process_data
                )
                
                if response.status_code != 200:
                    logger.error(f"Failed to process task: {response.text}")
                    retry_count += 1
                    time.sleep(3)
                    continue
                
                response = requests.get(
                    f"https://{server}/v1/download/{task_id}",
                    headers={"Authorization": f"Bearer {token}"},
                    stream=True
                )
                
                if response.status_code != 200:
                    logger.error(f"Failed to download file: {response.text}")
                    retry_count += 1
                    time.sleep(3)
                    continue
                    
                with open(pdf_path, 'wb') as f:
                    for chunk in response.iter_content(chunk_size=8192):
                        f.write(chunk)
                        
                if os.path.exists(pdf_path):
                    logger.info(f"Successfully created PDF: {pdf_path}")
                    return True
                else:
                    logger.error(f"PDF was not created at expected location: {pdf_path}")
                    retry_count += 1
                    time.sleep(3)
                    continue
                
            except Exception as e:
                logger.error(f"Error in API call (attempt {retry_count+1}): {e}")
                retry_count += 1
                time.sleep(3)
                
        logger.error(f"Failed to convert PDF after {max_retries} attempts")
        return False
            
    except Exception as e:
        logger.error(f"Error converting PPTX to PDF using ilovepdf: {e}")
        return False


def print_pdf(pdf_path):
    """Print PDF file to default printer on macOS"""
    try:
        logger.info(f"Printing PDF: {pdf_path}")
        result = subprocess.run(['lpr', pdf_path], check=True, capture_output=True)
        logger.info(f"Print job sent successfully for {pdf_path}")
        return True
    except subprocess.CalledProcessError as e:
        logger.error(f"Error printing PDF: {e}")
        logger.error(f"STDOUT: {e.stdout.decode('utf-8') if e.stdout else 'None'}")
        logger.error(f"STDERR: {e.stderr.decode('utf-8') if e.stderr else 'None'}")
        return False
    except Exception as e:
        logger.error(f"Unexpected error printing PDF: {e}")
        return False

def create_html_email(template_path, name):
    """Create HTML email with name replacing the placeholder"""
    try:
        with open(template_path, 'r') as file:
            html_content = file.read()
        
        html_content = html_content.replace('NAME_PLACEHOLDER', name)
        return html_content
    except Exception as e:
        logger.error(f"Error creating email content: {e}")
        return None

def send_email(recipient, subject, html_content, smtp_config, attachment_path=None):
    """Send email with certificate attachment using SMTP"""
    try:
        logger.info(f"Sending email to {recipient} using SMTP")
        
        # Check if SMTP config is provided
        if not smtp_config or not all(key in smtp_config for key in ["server", "port", "email", "password"]):
            logger.error("Missing SMTP configuration. Please add smtp settings to config.json")
            return False
        
        # Create message
        message = MIMEMultipart()
        message["From"] = smtp_config["email"]
        message["To"] = recipient
        message["Subject"] = subject
        
        # Attach HTML body
        message.attach(MIMEText(html_content, "html"))
        
        # Attach certificate if provided
        if attachment_path and os.path.exists(attachment_path):
            with open(attachment_path, "rb") as attachment:
                part = MIMEBase("application", "octet-stream")
                part.set_payload(attachment.read())
            
            encoders.encode_base64(part)
            part.add_header(
                "Content-Disposition",
                f"attachment; filename={os.path.basename(attachment_path)}",
            )
            message.attach(part)
        
        # Send email
        with smtplib.SMTP_SSL(smtp_config["server"], smtp_config["port"]) as server:
            server.login(smtp_config["email"], smtp_config["password"])
            server.send_message(message)
        logger.info(f"Email sent to {recipient}")
        return True
    except Exception as e:
        logger.error(f"Error sending email to {recipient}: {e}")
        return False
    
    #comment out above code if you dont want to send email
    return True

def generate_certificates():
    """Main function to generate and send certificates"""
    config = read_config()
    output_dir = config["output_dir"]
    Path(output_dir).mkdir(parents=True, exist_ok=True)
    users = read_csv(config["csv_path"])
    
    success_count = 0
    fail_count = 0
    
    for _, row in users.iterrows():
        try:
            name = row['name']
            email = row['email']
            safe_email = re.sub(r'[^\w\.-]', '_', email)
            pptx_path = os.path.join(output_dir, f"{safe_email}.pptx")
            pdf_path = os.path.join(output_dir, f"{safe_email}.pdf")
            
            logger.info(f"Generating certificate for {name} ({email})")
            
            if replace_placeholder_in_ppt(
                config["ppt_template_path"], 
                pptx_path, 
                "NAME_PLACEHOLDER", 
                name
            ):
                # Convert to PDF using ilovepdf (easier than handling PPTX directly on macos)
                # Use this method if you have a small number of certificates to generate
                # For larger numbers, consider using a windows machine with PowerPoint installed and COM automation using comtypes library
                logger.info(f"Converting certificate to PDF for {name} ({email})")
                if convert_pptx_to_pdf_using_ilovepdf(
                    pptx_path, 
                    pdf_path,
                    api_key=config.get("ilovepdf_public_key"),
                    api_secret=config.get("ilovepdf_secret_key")
                ):
                    # Print the PDF certificate
                    # print_success = print_pdf(pdf_path)
                    # if print_success:
                    #     logger.info(f"Certificate printed successfully for {name}")
                    # else:
                    #     logger.warning(f"Failed to print certificate for {name}")

                    html_content = create_html_email(config["html_template_path"], name)
                    
                    if html_content and send_email(
                        recipient=email,
                        subject=config.get("email_subject", "Your Certificate"),
                        html_content=html_content,
                        smtp_config=config.get("smtp"),
                        attachment_path=pdf_path
                    ):
                        success_count += 1
                    else:
                        fail_count += 1
                else:
                    logger.error(f"Failed to convert certificate to PDF for {name}")
                    fail_count += 1
            else:
                fail_count += 1
                
        except Exception as e:
            logger.error(f"Error processing certificate for {row.get('name', 'unknown')}: {e}")
            fail_count += 1
    
    logger.info(f"Certificate generation complete. Success: {success_count}, Failed: {fail_count}")

if __name__ == "__main__":
    generate_certificates()
